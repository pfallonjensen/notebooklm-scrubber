"""
Codia Visual Struct JSON → PPTX Converter

Converts Codia's structured JSON output into PowerPoint presentations
using python-pptx with high fidelity positioning and styling.
"""

import json
import io
import requests
from pathlib import Path
from typing import Optional, Tuple, List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# Import OCR corrector
try:
    from ocr_corrector import correct_text, is_garbage_text, apply_regex_fixes
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


# Standard slide dimensions (16:9 widescreen)
SLIDE_WIDTH_INCHES = 10
SLIDE_HEIGHT_INCHES = 7.5

# Codia's standard output dimensions
CODIA_BASE_WIDTH = 2867
CODIA_BASE_HEIGHT = 1600

# Font mapping: Codia detected fonts → safe PPTX fonts
FONT_MAP = {
    'Rethink Sans': 'Arial',
    'Inter': 'Arial',
    'Manrope': 'Arial',
    'DM Sans': 'Arial',
    'Gabarito': 'Arial',
    'Onest': 'Arial',
    'Epilogue': 'Arial',
    'Figtree': 'Arial',
    'Reddit Sans': 'Arial',
    'Space Grotesk': 'Arial',
    'Hedvig Letters Sans': 'Arial',
    'Cabin': 'Arial',
    'Archivo': 'Arial',
    'Roboto Flex': 'Arial',
    'Geist': 'Arial',
    'Plus Jakarta Sans': 'Arial',
    'Quicksand': 'Arial',
    'Urbanist': 'Arial',
    'Work Sans': 'Arial',
    'Barlow': 'Arial',
    'Poppins': 'Arial',
    'News Cycle': 'Arial',
    # Add more as discovered
}


class CodiaConverter:
    """Converts Codia Visual Struct JSON to PPTX."""

    def __init__(self, codia_json: dict, ocr_correction: bool = True, use_llm: bool = False):
        """
        Initialize converter with Codia JSON data.

        Args:
            codia_json: Parsed JSON from Codia Visual Struct API
            ocr_correction: Enable OCR error correction (regex-based)
            use_llm: Enable LLM-based OCR correction (requires anthropic package)
        """
        self.data = codia_json
        self.config = codia_json.get('visual_struct', {}).get('configuration', {})
        self.ocr_correction = ocr_correction and OCR_AVAILABLE
        self.use_llm = use_llm

        # Get actual dimensions from Codia config
        self.codia_width = self.config.get('baseWidth', CODIA_BASE_WIDTH)
        self.codia_height = CODIA_BASE_HEIGHT  # Height not in config, using default

        # Calculate scale factors
        self.scale_x = Inches(SLIDE_WIDTH_INCHES) / self.codia_width
        self.scale_y = Inches(SLIDE_HEIGHT_INCHES) / self.codia_height

        # Track downloaded images to avoid re-downloading
        self._image_cache: Dict[str, bytes] = {}

        # Collect all text for OCR correction (batching)
        self._text_elements: List[Dict] = []

    def convert(self, output_path: str) -> str:
        """
        Convert Codia JSON to PPTX file.

        Args:
            output_path: Path to save the PPTX file

        Returns:
            Path to the created PPTX file
        """
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

        # Add a blank slide (layout index 6 is typically blank)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Get root element
        root = self.data.get('visual_struct', {}).get('visualElement', {})

        # Process background color first
        self._apply_background(slide, root)

        # Process all elements recursively
        # Images first (background layer), then text (foreground)
        self._collect_elements(root)

        # Add images first
        for element in self._get_image_elements(root):
            self._add_image(slide, element)

        # Add text elements on top
        for element in self._get_text_elements(root):
            self._add_text(slide, element)

        prs.save(output_path)
        return output_path

    def _apply_background(self, slide, root: dict):
        """Apply background color from root element."""
        style = root.get('styleConfig', {})
        bg_spec = style.get('backgroundSpec', {})

        if bg_spec.get('type') == 'COLOR':
            rgb = bg_spec.get('backgroundColor', {}).get('rgb', [255, 255, 255])
            if rgb and len(rgb) >= 3:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])

    def _collect_elements(self, element: dict, parent_offset: Tuple[int, int] = (0, 0)) -> List[dict]:
        """Recursively collect all elements with absolute positions."""
        elements = []

        # Get this element's position
        layout = element.get('layoutConfig', {})
        abs_attrs = layout.get('absoluteAttrs', {})

        # Use orginCoord for original positioning
        coord = abs_attrs.get('orginCoord', abs_attrs.get('coord', [0, 0]))
        if coord and len(coord) >= 2:
            abs_x = coord[0]
            abs_y = coord[1]
        else:
            abs_x, abs_y = 0, 0

        # Add this element's position
        element['_abs_x'] = abs_x
        element['_abs_y'] = abs_y

        elements.append(element)

        # Recurse into children
        for child in element.get('childElements', []):
            elements.extend(self._collect_elements(child, (abs_x, abs_y)))

        return elements

    def _get_image_elements(self, root: dict) -> List[dict]:
        """Get all image elements, sorted by z-order (depth-first)."""
        elements = self._collect_elements(root)
        return [e for e in elements if e.get('elementType') == 'Image']

    def _get_text_elements(self, root: dict) -> List[dict]:
        """Get all text elements, sorted by z-order (depth-first)."""
        elements = self._collect_elements(root)
        return [e for e in elements if e.get('elementType') == 'Text']

    def _add_image(self, slide, element: dict):
        """Add an image element to the slide."""
        content = element.get('contentData', {})
        image_url = content.get('imageSource', '')

        if not image_url:
            return

        # Get position and size
        x = element.get('_abs_x', 0)
        y = element.get('_abs_y', 0)

        style = element.get('styleConfig', {})
        width = style.get('widthSpec', {}).get('value', 100)
        height = style.get('heightSpec', {}).get('value', 100)

        # Convert to EMUs
        left = int(x * self.scale_x)
        top = int(y * self.scale_y)
        img_width = int(width * self.scale_x)
        img_height = int(height * self.scale_y)

        # Download image
        try:
            if image_url in self._image_cache:
                img_bytes = self._image_cache[image_url]
            else:
                response = requests.get(image_url, timeout=30)
                response.raise_for_status()
                img_bytes = response.content
                self._image_cache[image_url] = img_bytes

            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, left, top, img_width, img_height)
        except Exception as e:
            print(f"Warning: Failed to download image {image_url}: {e}")

    def _add_text(self, slide, element: dict):
        """Add a text element to the slide."""
        content = element.get('contentData', {})
        text = content.get('textValue', '')

        if not text or text.strip() == '':
            return

        # Apply OCR correction if enabled
        if self.ocr_correction:
            if is_garbage_text(text):
                return
            text, fixes = apply_regex_fixes(text)
            if fixes:
                print(f"  OCR fixes: {fixes}")
        elif self._is_garbage_text(text):
            return

        # Get position and size
        x = element.get('_abs_x', 0)
        y = element.get('_abs_y', 0)

        style = element.get('styleConfig', {})
        width = style.get('widthSpec', {}).get('value', 200)
        height = style.get('heightSpec', {}).get('value', 50)

        # Convert to EMUs
        left = int(x * self.scale_x)
        top = int(y * self.scale_y)
        box_width = int(width * self.scale_x)
        box_height = int(height * self.scale_y)

        # Create text box
        textbox = slide.shapes.add_textbox(left, top, box_width, box_height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        # Get text styling
        text_config = style.get('textConfig', {})
        font_family = text_config.get('fontFamily', 'Arial')
        font_size = text_config.get('fontSize', 12)
        font_style = text_config.get('fontStyle', 'regular')

        # Map font to safe alternative
        font_family = FONT_MAP.get(font_family, font_family)

        # Get text color
        text_color = style.get('textColor', {})
        rgb_values = text_color.get('rgbValues', [0, 0, 0])

        # Get alignment
        text_align = text_config.get('textAlign', ['LEFT', 'CENTER'])
        h_align = text_align[0] if text_align else 'LEFT'

        # Set paragraph properties
        p = tf.paragraphs[0]
        p.text = text

        # Set alignment
        if h_align == 'CENTER':
            p.alignment = PP_ALIGN.CENTER
        elif h_align == 'RIGHT':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        # Set font properties
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_family
        run.font.size = Pt(font_size)

        # Set bold/italic based on font_style
        if 'bold' in font_style.lower() or 'extra_bold' in font_style.lower():
            run.font.bold = True
        if 'italic' in font_style.lower():
            run.font.italic = True

        # Set color
        if rgb_values and len(rgb_values) >= 3:
            run.font.color.rgb = RGBColor(rgb_values[0], rgb_values[1], rgb_values[2])

    def _is_garbage_text(self, text: str) -> bool:
        """Check if text is likely OCR garbage to be skipped."""
        garbage_patterns = [
            'i\nS',  # Common icon misread
            '→0',   # Arrow symbols misread
        ]
        return text.strip() in garbage_patterns

    def get_all_text(self) -> List[str]:
        """
        Get all text content for OCR correction batching.

        Returns:
            List of text strings from all text elements
        """
        root = self.data.get('visual_struct', {}).get('visualElement', {})
        text_elements = self._get_text_elements(root)
        return [
            e.get('contentData', {}).get('textValue', '')
            for e in text_elements
            if e.get('contentData', {}).get('textValue', '')
        ]


def convert_codia_to_pptx(json_path: str, output_path: str) -> str:
    """
    Convenience function to convert a Codia JSON file to PPTX.

    Args:
        json_path: Path to Codia JSON file
        output_path: Path for output PPTX file

    Returns:
        Path to created PPTX file
    """
    with open(json_path, 'r') as f:
        codia_json = json.load(f)

    converter = CodiaConverter(codia_json)
    return converter.convert(output_path)


def convert_multiple_pages(json_paths: List[str], output_path: str) -> str:
    """
    Convert multiple Codia JSON files (one per page) into a single PPTX.

    Args:
        json_paths: List of paths to Codia JSON files
        output_path: Path for output PPTX file

    Returns:
        Path to created PPTX file
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    for json_path in json_paths:
        with open(json_path, 'r') as f:
            codia_json = json.load(f)

        converter = CodiaConverter(codia_json)

        # Add slide for this page
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        root = codia_json.get('visual_struct', {}).get('visualElement', {})

        # Apply background
        converter._apply_background(slide, root)

        # Add images first, then text
        for element in converter._get_image_elements(root):
            converter._add_image(slide, element)

        for element in converter._get_text_elements(root):
            converter._add_text(slide, element)

    prs.save(output_path)
    return output_path


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python codia_to_pptx.py <input.json> <output.pptx>")
        print("       python codia_to_pptx.py --multi <output.pptx> <page1.json> <page2.json> ...")
        sys.exit(1)

    if sys.argv[1] == '--multi':
        output = sys.argv[2]
        inputs = sys.argv[3:]
        result = convert_multiple_pages(inputs, output)
    else:
        input_path = sys.argv[1]
        output_path = sys.argv[2]
        result = convert_codia_to_pptx(input_path, output_path)

    print(f"Created: {result}")
